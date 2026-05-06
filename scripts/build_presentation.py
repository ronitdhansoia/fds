"""Build a black-themed Keynote-style 16:9 .pptx for the MigrantMoney FDS talk.

Run:
    uv run --with python-pptx --with Pillow python scripts/build_presentation.py

Produces: report/MigrantMoney_Presentation.pptx
"""
from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Design tokens (mirror the dashboard palette)
# ---------------------------------------------------------------------------
BG = RGBColor(0x0A, 0x0A, 0x0A)
SURFACE = RGBColor(0x14, 0x14, 0x14)
BORDER = RGBColor(0x24, 0x24, 0x24)
TEXT = RGBColor(0xF2, 0xF2, 0xF2)
TEXT_2 = RGBColor(0xA8, 0xA8, 0xA8)
TEXT_3 = RGBColor(0x6B, 0x6B, 0x6B)
ACCENT = RGBColor(0xE8, 0x86, 0x18)
ACCENT_2 = RGBColor(0x86, 0xC2, 0x2C)
WHITE_CARD = RGBColor(0xFA, 0xFA, 0xFA)

DISPLAY = "Helvetica Neue"
BODY = "Helvetica"
MONO = "Menlo"

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "report" / "figures"
META_PATH = ROOT / "data" / "outputs" / "meta.json"
OUT = ROOT / "report" / "MigrantMoney_Presentation.pptx"


# ---------------------------------------------------------------------------
# Presentation skeleton
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    """Create a fresh black-background slide."""
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return s


def text(
    slide,
    x,
    y,
    w,
    h,
    content,
    *,
    size=18,
    color=TEXT,
    font=BODY,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=None,
    spacing=None,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    p.text = content
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def overline(slide, x, y, content, *, color=TEXT_3):
    return text(slide, x, y, Inches(10), Inches(0.3), content,
                size=10, color=color, font=MONO)


def divider(slide, x, y, length, color=BORDER, weight=0.6):
    """Hairline rule."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + length, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def page_chrome(slide, page_no, total, section):
    """Top-of-slide rule + tiny page indicator + section overline."""
    overline(slide, Inches(0.85), Inches(0.55), section)
    text(slide, Inches(11.5), Inches(0.55), Inches(1.2), Inches(0.3),
         f"{page_no:02d} / {total:02d}",
         size=10, color=TEXT_3, font=MONO, align=PP_ALIGN.RIGHT)
    divider(slide, Inches(0.85), Inches(0.95), Inches(11.65))


def figure(slide, fig_name, *, top=Inches(2.3), max_w=Inches(10.5),
           max_h=Inches(4.6), card=True):
    """Place a figure centered with optional white card behind."""
    p = FIG / fig_name
    with Image.open(p) as im:
        iw, ih = im.size
    target_w = max_w
    target_h = int(target_w * ih / iw)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * iw / ih)
    left = (W - target_w) // 2
    if card:
        pad = Inches(0.18)
        c = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left - pad, top - pad,
            target_w + 2 * pad, target_h + 2 * pad,
        )
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE_CARD
        c.line.fill.background()
        c.adjustments[0] = 0.012
    slide.shapes.add_picture(str(p), left, top, width=target_w, height=target_h)


# Load meta for headline numbers
META = json.loads(META_PATH.read_text())
g = META["global_savings"]
TOTAL_VOLUME_B = g["total_corridor_volume_usd"] / 1e9
SAVINGS_B = g["total_savings_usd_annual_current"] / 1e9
SAVINGS_PCT = g["implied_avg_savings_pct_current"]
N_POSITIVE = g["n_corridors_with_positive_savings"]
N_WITH_VOL = g["n_corridors_with_volume"]
N_CORRIDORS = META["n_corridors"]
N_PROVIDERS = META["n_providers"]
N_ROWS = META["n_rows"]
N_QUARTERS = META["n_quarters"]
PERIOD_FROM = META["panel_first_period"].replace("_", " ")
PERIOD_TO = META["panel_last_period"].replace("_", " ")


# Total slide count for page numbers
TOTAL = 18

# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s = add_slide()
overline(s, Inches(0.85), Inches(0.55),
         "FUNDAMENTALS OF DATA SCIENCE  ·  BITS PILANI DUBAI  ·  2026")
divider(s, Inches(0.85), Inches(0.95), Inches(11.65))

# Hero
text(s, Inches(0.85), Inches(2.45), Inches(12.0), Inches(2.4),
     "MigrantMoney",
     size=110, color=TEXT, font=DISPLAY, bold=True)

# Subhead
text(s, Inches(0.85), Inches(4.55), Inches(11.5), Inches(0.6),
     "A True Cost Index and stablecoin counterfactual for the global remittance market.",
     size=22, color=TEXT_2, font=BODY)

# Footer
divider(s, Inches(0.85), Inches(6.6), Inches(11.65))
text(s, Inches(0.85), Inches(6.78), Inches(6), Inches(0.4),
     "RONIT DHANSOIA",
     size=10, color=TEXT_2, font=MONO)
text(s, Inches(7.0), Inches(6.78), Inches(5.5), Inches(0.4),
     "TERM PROJECT  ·  MAY 2026",
     size=10, color=TEXT_3, font=MONO, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 2 — The question (hero number)
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 2, TOTAL, "01  ·  THE QUESTION")

text(s, Inches(0.85), Inches(2.0), Inches(12), Inches(2.5),
     f"${SAVINGS_B:.2f} billion",
     size=120, color=TEXT, font=DISPLAY, bold=True)

text(s, Inches(0.85), Inches(4.4), Inches(12), Inches(0.6),
     "What if the world's most expensive remittance corridors",
     size=28, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(4.95), Inches(12), Inches(0.6),
     "ran on stablecoin rails? An empirical answer, per corridor.",
     size=28, color=TEXT_2, font=BODY)

text(s, Inches(0.85), Inches(6.6), Inches(12), Inches(0.4),
     f"World Bank panel, {PERIOD_FROM} → {PERIOD_TO}  ·  {N_CORRIDORS} corridors  ·  {N_ROWS:,} observations.",
     size=11, color=TEXT_3, font=MONO)


# ---------------------------------------------------------------------------
# Slide 3 — The gap
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 3, TOTAL, "02  ·  THE GAP")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Three signals, never combined.",
     size=44, color=TEXT, font=DISPLAY, bold=True)

text(s, Inches(0.85), Inches(2.7), Inches(11), Inches(1.4),
     "Public reports treat fee and FX margin as separate columns. Speed isn't priced at all. No published metric unifies the three; we engineer one.",
     size=18, color=TEXT_2, font=BODY, line_spacing=1.4)

# Three columns: fee, fx margin, speed penalty
col_w = Inches(3.7)
col_y = Inches(4.6)
col_pad = Inches(0.35)
xs = [Inches(0.85), Inches(0.85) + col_w + col_pad, Inches(0.85) + 2 * (col_w + col_pad)]
labels = ["FEE %", "FX MARGIN %", "SPEED PENALTY"]
descs = [
    "What the receipt shows. Always advertised.",
    "Spread between the rate offered and the interbank mid. Rarely shown.",
    "Cost to the receiving household of waiting two or three days.",
]
for x, label, desc in zip(xs, labels, descs):
    text(s, x, col_y, col_w, Inches(0.4), label,
         size=10, color=ACCENT, font=MONO, bold=True)
    divider(s, x, col_y + Inches(0.45), col_w, color=BORDER, weight=0.4)
    text(s, x, col_y + Inches(0.62), col_w, Inches(2),
         desc, size=15, color=TEXT_2, font=BODY, line_spacing=1.4)


# ---------------------------------------------------------------------------
# Slide 4 — Objectives
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 4, TOTAL, "03  ·  OBJECTIVES")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Two outputs.",
     size=52, color=TEXT, font=DISPLAY, bold=True)

# Objective 1
text(s, Inches(0.85), Inches(3.1), Inches(0.5), Inches(0.5),
     "01", size=26, color=TEXT_3, font=MONO)
text(s, Inches(1.7), Inches(3.05), Inches(11), Inches(0.7),
     "True Cost Index.",
     size=30, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(1.7), Inches(3.75), Inches(11), Inches(1.2),
     "A unified per-corridor cost number that combines fee, FX margin, and a speed penalty.",
     size=17, color=TEXT_2, font=BODY, line_spacing=1.4)

# Objective 2
text(s, Inches(0.85), Inches(5.3), Inches(0.5), Inches(0.5),
     "02", size=26, color=TEXT_3, font=MONO)
text(s, Inches(1.7), Inches(5.25), Inches(11), Inches(0.7),
     "Stablecoin counterfactual.",
     size=30, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(1.7), Inches(5.95), Inches(11), Inches(1.2),
     "Per-corridor savings if those flows ran on USDC / USDT rails. The first published estimate at this granularity.",
     size=17, color=TEXT_2, font=BODY, line_spacing=1.4)


# ---------------------------------------------------------------------------
# Slide 5 — Dataset
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 5, TOTAL, "04  ·  DATASET")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "World Bank Remittance Prices Worldwide.",
     size=36, color=TEXT, font=DISPLAY, bold=True)

text(s, Inches(0.85), Inches(2.55), Inches(11), Inches(0.6),
     f"Quarterly panel of corridor prices. {PERIOD_FROM} → {PERIOD_TO}.",
     size=18, color=TEXT_2, font=BODY)

# Number grid 2x2
def stat_block(slide, x, y, label, value, sub):
    text(slide, x, y, Inches(5.5), Inches(0.4), label,
         size=10, color=TEXT_3, font=MONO)
    text(slide, x, y + Inches(0.45), Inches(5.5), Inches(1.0), value,
         size=64, color=TEXT, font=DISPLAY, bold=True)
    text(slide, x, y + Inches(1.55), Inches(5.5), Inches(0.4), sub,
         size=12, color=TEXT_2, font=BODY)

stat_block(s, Inches(0.85), Inches(3.7), "CORRIDORS", f"{N_CORRIDORS}", "country pair × send amount")
stat_block(s, Inches(7.0), Inches(3.7), "PROVIDERS", f"{N_PROVIDERS}", "MTOs, banks, fintechs, mobile money")
stat_block(s, Inches(0.85), Inches(5.55), "QUARTERS", f"{N_QUARTERS}", f"{PERIOD_FROM} through {PERIOD_TO}")
stat_block(s, Inches(7.0), Inches(5.55), "ROWS", f"{N_ROWS:,}", "after schema sniff and clean")


# ---------------------------------------------------------------------------
# Slide 6 — Pipeline (8 stages)
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 6, TOTAL, "05  ·  PIPELINE")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Pipeline.",
     size=52, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(2.6), Inches(11.5), Inches(0.5),
     "Eight stages from raw RPW xlsx to a static dashboard. Reproducible end-to-end.",
     size=16, color=TEXT_2, font=BODY)

stages = [
    ("01", "INGEST",      "Download RPW xlsx + KNOMAD bilateral matrix"),
    ("02", "PREPROCESS",  "Schema sniff  ·  cc1 / cc2 melt  ·  feature derivation"),
    ("03", "TCI",         "Corridor × period × amount aggregation"),
    ("04", "STABLECOIN",  "Counterfactual cost model + savings"),
    ("05", "REGRESSION",  "Two-way fixed effects  ·  clustered SE"),
    ("06", "AGGREGATE",   "Diaspora burden + corridor rankings"),
    ("07", "EXPORT",      "Round  ·  null-safe JSON for the dashboard"),
    ("08", "FIGURES",     "Plotly  ·  300 DPI PNG for the report"),
]

# Render two columns of 4 each
y_start = Inches(3.45)
row_h = Inches(0.62)
for i, (num, name, desc) in enumerate(stages):
    col = i // 4
    row = i % 4
    x = Inches(0.85) + col * Inches(6.2)
    y = y_start + row * row_h
    text(s, x, y, Inches(0.5), Inches(0.4), num,
         size=11, color=TEXT_3, font=MONO)
    text(s, x + Inches(0.55), y, Inches(2.2), Inches(0.4), name,
         size=12, color=TEXT, font=MONO, bold=True)
    text(s, x + Inches(2.5), y, Inches(3.6), Inches(0.4), desc,
         size=12, color=TEXT_2, font=BODY)

text(s, Inches(0.85), Inches(6.55), Inches(11.5), Inches(0.4),
     "RUN  ·  uv run python scripts/run_all.py    ~ 90 s end-to-end on cached data.",
     size=10, color=TEXT_3, font=MONO)


# ---------------------------------------------------------------------------
# Slide 7 — Data engineering
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 7, TOTAL, "06  ·  DATA ENGINEERING")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Data engineering.",
     size=44, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(2.45), Inches(11.5), Inches(0.5),
     "Six transformations turn raw RPW xlsx into a clean panel ready for analysis.",
     size=16, color=TEXT_2, font=BODY)

steps = [
    ("01", "SCHEMA SNIFF", "Match column names against canonical CANDIDATE_COLUMNS  ·  RPW renames quarterly"),
    ("02", "MELT",         "Per-amount columns (cc1 = USD 200, cc2 = USD 500)  →  long form"),
    ("03", "DERIVE FEE",   "fee_pct = total_cost_pct − fx_margin_pct  ·  clipped at zero"),
    ("04", "MAP SPEED",    'RPW "speed actual" categorical  →  days_to_arrive scalar  →  speed penalty'),
    ("05", "FILTER",       "Drop rows missing any TCI input  ·  log dropped count for QA"),
    ("06", "NORMALIZE",    "Standardize firm_type to {Bank, MTO, MobileMoney, PostOffice, Fintech}"),
]

y_start = Inches(3.35)
row_h = Inches(0.55)
for i, (num, name, desc) in enumerate(steps):
    y = y_start + i * row_h
    text(s, Inches(0.85), y, Inches(0.5), Inches(0.4), num,
         size=11, color=TEXT_3, font=MONO)
    text(s, Inches(1.4), y, Inches(2.6), Inches(0.4), name,
         size=12, color=TEXT, font=MONO, bold=True)
    text(s, Inches(3.7), y, Inches(8.5), Inches(0.4), desc,
         size=12, color=TEXT_2, font=BODY)

text(s, Inches(0.85), Inches(6.85), Inches(11.5), Inches(0.4),
     f"OUTPUT  ·  {N_ROWS:,} clean rows  ·  Parquet for downstream stages.",
     size=10, color=TEXT_3, font=MONO)


# ---------------------------------------------------------------------------
# Slide 8 — TCI formula
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 8, TOTAL, "07  ·  METHODOLOGY  ·  TRUE COST INDEX")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "True Cost Index.",
     size=44, color=TEXT, font=DISPLAY, bold=True)

# Formula
text(s, Inches(0.85), Inches(3.0), Inches(12), Inches(1.2),
     "TCI  =  fee%  +  fxMargin%  +  κ · max(0, days−1)",
     size=26, color=TEXT, font=MONO)

text(s, Inches(0.85), Inches(4.4), Inches(11.5), Inches(0.4),
     "WHERE",
     size=10, color=TEXT_3, font=MONO)
text(s, Inches(0.85), Inches(4.78), Inches(11.5), Inches(0.5),
     "κ = 0.10 % per day past same-day  ·  cost-of-capital proxy for the receiving household",
     size=15, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(5.3), Inches(11.5), Inches(0.5),
     'days  ·  RPW "speed actual" mapped to (under 1h, same day) → 0,  next day → 1,  2d → 2,  3-5d → 4',
     size=15, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(5.82), Inches(11.5), Inches(0.5),
     "Corridor-level TCI is the unweighted mean across providers; median reported as a robustness check.",
     size=15, color=TEXT_2, font=BODY)


# ---------------------------------------------------------------------------
# Slide 9 — Stablecoin model
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 9, TOTAL, "08  ·  METHODOLOGY  ·  STABLECOIN")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Stablecoin counterfactual.",
     size=44, color=TEXT, font=DISPLAY, bold=True)

text(s, Inches(0.85), Inches(2.85), Inches(12), Inches(1.0),
     "SC = onramp(s)  +  offramp(d)  +  (gas / A × 100)  +  fxSpread(d)",
     size=22, color=TEXT, font=MONO)
text(s, Inches(0.85), Inches(3.55), Inches(12), Inches(1.0),
     "savings = max(0,  TCI  −  SC)  ·  V(s,d)",
     size=22, color=TEXT, font=MONO)

text(s, Inches(0.85), Inches(4.55), Inches(11.5), Inches(0.4),
     "DEFAULTS",
     size=10, color=TEXT_3, font=MONO)
text(s, Inches(0.85), Inches(4.93), Inches(11.5), Inches(0.5),
     "On-ramp  ·  1.0 % developed  /  1.5 % default  /  2.5 % low-banked",
     size=14, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(5.38), Inches(11.5), Inches(0.5),
     "Off-ramp  ·  1.0 % top P2P  /  2.5 % default  /  4.0 % thin liquidity",
     size=14, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(5.83), Inches(11.5), Inches(0.5),
     "Gas $0.50 (L2 / Solana / Tron USDT)  ·  fxSpread 0.5 % deep / 1.5 % default",
     size=14, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(6.28), Inches(11.5), Inches(0.5),
     "Volume V(s,d) is the KNOMAD bilateral remittance estimate, 2021 (latest published).",
     size=14, color=TEXT_3, font=BODY, italic=True)


# ---------------------------------------------------------------------------
# Slide 10 — Regression spec
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 10, TOTAL, "09  ·  METHODOLOGY  ·  REGRESSION")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Two-way fixed-effects panel.",
     size=40, color=TEXT, font=DISPLAY, bold=True)

text(s, Inches(0.85), Inches(2.85), Inches(12), Inches(1.0),
     "TCI_ipq  =  β₀  +  Σ β_k · 1{firmType = k}  +  α_corridor  +  γ_quarter  +  ε",
     size=22, color=TEXT, font=MONO)

text(s, Inches(0.85), Inches(4.2), Inches(11.5), Inches(0.5),
     "Entity FE  ·  corridor (368 levels)  ·  absorbs all corridor-invariant heterogeneity.",
     size=18, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(4.85), Inches(11.5), Inches(0.5),
     "Time FE  ·  quarter (36 levels)  ·  absorbs aggregate time trends.",
     size=18, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(5.5), Inches(11.5), Inches(0.5),
     "Reference category  ·  MTO (largest cell)  ·  cluster-robust SE at the corridor level.",
     size=18, color=TEXT_2, font=BODY)
text(s, Inches(0.85), Inches(6.15), Inches(11.5), Inches(0.5),
     "Implementation  ·  linearmodels.PanelOLS  ·  fit separately for the USD 200 and USD 500 buckets.",
     size=18, color=TEXT_2, font=BODY)


# ---------------------------------------------------------------------------
# Slide 11 — Architecture diagram
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 11, TOTAL, "10  ·  ARCHITECTURE")

text(s, Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.7),
     "From raw RPW to a static dashboard.",
     size=28, color=TEXT, font=DISPLAY, bold=True)

figure(s, "architecture.png", top=Inches(2.15), max_w=Inches(10.5), max_h=Inches(4.9), card=True)


# ---------------------------------------------------------------------------
# Slide 12 — Top corridors
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 12, TOTAL, "11  ·  RESULTS  ·  COSTS")

text(s, Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.7),
     "Top 20 most expensive corridors.",
     size=28, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.4),
     "Sub-Saharan Africa and small-island Pacific dominate. Latin America is the cheapest.",
     size=14, color=TEXT_2, font=BODY)

figure(s, "fig01_top20_corridors.png", top=Inches(2.45), max_w=Inches(10.0), max_h=Inches(4.7))


# ---------------------------------------------------------------------------
# Slide 13 — Regression forest plot
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 13, TOTAL, "12  ·  RESULTS  ·  OPERATORS")

text(s, Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.7),
     "Banks charge ≈ 4.5 pp more than MTOs.",
     size=28, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.4),
     "After controlling for corridor and quarter. Mobile-money operators undercut MTOs by ≈ 1.5 pp.",
     size=14, color=TEXT_2, font=BODY)

figure(s, "fig03_operator_forest.png", top=Inches(2.45), max_w=Inches(10.0), max_h=Inches(4.7))


# ---------------------------------------------------------------------------
# Slide 14 — Stablecoin scatter
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 14, TOTAL, "13  ·  RESULTS  ·  STABLECOIN")

text(s, Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.7),
     "Where stablecoin rails save the most.",
     size=28, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.4),
     "Each dot is a corridor. X-axis: corridor volume (log). Y-axis: percentage savings on USD 200.",
     size=14, color=TEXT_2, font=BODY)

figure(s, "fig04_stablecoin_scatter.png", top=Inches(2.45), max_w=Inches(10.0), max_h=Inches(4.7))


# ---------------------------------------------------------------------------
# Slide 15 — Diaspora burden
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 15, TOTAL, "14  ·  RESULTS  ·  DIASPORA")

text(s, Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.7),
     "Annual burden by sending country.",
     size=28, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.4),
     "How many dollars each diaspora pays per year on the corridors out of their host country.",
     size=14, color=TEXT_2, font=BODY)

figure(s, "fig05_diaspora_burden.png", top=Inches(2.45), max_w=Inches(10.0), max_h=Inches(4.7))


# ---------------------------------------------------------------------------
# Slide 16 — Robustness
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 16, TOTAL, "15  ·  ROBUSTNESS")

text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(0.9),
     "Does the headline survive scrutiny?",
     size=40, color=TEXT, font=DISPLAY, bold=True)
text(s, Inches(0.85), Inches(2.55), Inches(11.5), Inches(0.5),
     "Four checks. Every assumption is exposed on the dashboard.",
     size=16, color=TEXT_2, font=BODY)

checks = [
    ("01", "κ SENSITIVITY",
     "Reviewer-supplied κ ∈ [0.05, 0.20]  ·  corridor ranking is largely invariant in the plausible range."),
    ("02", "PROVIDER STRESS TEST",
     "Drop top + bottom 5 % of provider observations per corridor  ·  global savings move by < 3 %."),
    ("03", "AGGREGATION CHOICE",
     "Mean and median TCI reported alongside  ·  divergence < 0.5 pp on 80 % of corridors."),
    ("04", "ASSUMPTION FLEX",
     "Optimistic stablecoin defaults reach the $30 – 50 B / yr ballpark cited by BIS Bulletin 87."),
]

y_start = Inches(3.45)
row_h = Inches(0.78)
for i, (num, name, desc) in enumerate(checks):
    y = y_start + i * row_h
    text(s, Inches(0.85), y, Inches(0.5), Inches(0.4), num,
         size=11, color=TEXT_3, font=MONO)
    text(s, Inches(1.4), y, Inches(3.5), Inches(0.4), name,
         size=13, color=TEXT, font=MONO, bold=True)
    text(s, Inches(1.4), y + Inches(0.32), Inches(11), Inches(0.4),
         desc, size=13, color=TEXT_2, font=BODY)


# ---------------------------------------------------------------------------
# Slide 17 — Limitations
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 17, TOTAL, "16  ·  LIMITATIONS")

text(s, Inches(0.85), Inches(1.5), Inches(11.5), Inches(0.9),
     "What this can't tell you.",
     size=40, color=TEXT, font=DISPLAY, bold=True)

items = [
    ("Provider weighting is uniform.", "RPW does not publish market shares. A bank counts as much as a fintech."),
    ("κ is calibrated, not measured.", "Defensible at 0.10 % / day, but the slider on the dashboard exists for a reason."),
    ("Stablecoin off-ramp is a floor.", "KYC delays, P2P trade risk, and de-risking are not in the cost number."),
    ("BRM volumes are 2021.", "Latest published. Headline USD figure is anchored to that year."),
    ("No causal claim.", "Operator-class effect after FE; selection into firm type is endogenous."),
]

y = Inches(2.7)
for label, body in items:
    text(s, Inches(0.85), y, Inches(11), Inches(0.4),
         label, size=15, color=TEXT, font=BODY, bold=True)
    text(s, Inches(0.85), y + Inches(0.32), Inches(11), Inches(0.4),
         body, size=13, color=TEXT_2, font=BODY)
    y += Inches(0.78)


# ---------------------------------------------------------------------------
# Slide 18 — Demo / Q&A
# ---------------------------------------------------------------------------
s = add_slide()
page_chrome(s, 18, TOTAL, "17  ·  DEMO  ·  Q&A")

text(s, Inches(0.85), Inches(1.7), Inches(11.5), Inches(1.4),
     "Live demo.",
     size=88, color=TEXT, font=DISPLAY, bold=True)

text(s, Inches(0.85), Inches(3.6), Inches(11.5), Inches(0.6),
     "migrantmoney.vercel.app",
     size=28, color=ACCENT, font=MONO)
text(s, Inches(0.85), Inches(4.25), Inches(11.5), Inches(0.6),
     "github.com/ronitdhansoia/fds",
     size=18, color=TEXT_2, font=MONO)

text(s, Inches(0.85), Inches(6.5), Inches(11.5), Inches(0.4),
     "Ronit Dhansoia  ·  f20220168@dubai.bits-pilani.ac.in",
     size=11, color=TEXT_3, font=MONO)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)")
